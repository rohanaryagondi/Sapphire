"""Build the 6/7 follow-up deck from the OFFICIAL Quiver template.

Per the (2026-06-07-updated) /quiver-branding skill — references/slides.md gives
exact XML for every chrome element. The previous build missed:
  - the body-border rectangle (REQUIRED on every content slide)
  - chip adj=14493 (we used 50000 → fully-circular lozenge)
  - chip <a:noFill/> (we used white fill as a workaround)
  - title as a PLAIN TEXT BOX, not <p:ph type="title"/>
  - the slide-number placeholder per slide
This rewrite uses direct XML injection for the chrome and python-pptx for
content (tables/images/text). Body border is added FIRST in z-order so it
sits behind content.

Layouts used:
  0  Title Slide      — opening (uses template's own gradient)
  2  Blank White Layout — content slides; we delete its placeholders + inject
                          our own chrome (chip/title/body border/sldnum)
 12  Two boxes w/subtitle — for the image+text split (BBBP slide)
 42  Closing slide    — closing (template's gradient)

Run: /opt/anaconda3/envs/mammal/bin/python build_from_template.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

ROOT = Path(__file__).resolve().parent
SKILL_ASSETS = Path("/Users/rohanaryagondi/.claude/skills/quiver-branding/assets")
TEMPLATE = SKILL_ASSETS / "slideTemplateQuiver_2025.pptx"

OUTPUT = ROOT / "MAMMAL_followup_6-7.pptx"
FIG_CURVE = str(ROOT / "fig_datafit_curve.png")
FIG_BBBP = str(ROOT / "fig_bbbp_physchem.png")

# ---- Brand palette ----
PURPLE = RGBColor(0x9C, 0x02, 0xFA)
DARKP  = RGBColor(0x73, 0x00, 0xBF)
RED    = RGBColor(0xF5, 0x1D, 0x3F)
INK    = RGBColor(0x1A, 0x1A, 0x24)
GREY   = RGBColor(0x55, 0x55, 0x6A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FONT   = "Poppins"

CHIP = "MAMMAL"

# ---- Body-border bounds (per slides.md) ----
BB_X, BB_Y = 150876, 630936
BB_W, BB_H = 8833104, 4080510
# Content area inside the body border (with ~120000 EMU padding):
CONTENT_X = BB_X + 120000          # 270876
CONTENT_Y = BB_Y + 120000          # 750936
CONTENT_W = BB_W - 240000          # 8593104
CONTENT_H = BB_H - 240000          # 3840510

# Namespace declarations used by all injected XML fragments.
NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}
P_NS = NSMAP["p"]
A_NS = NSMAP["a"]


def _xml(fragment):
    """Parse an XML fragment with a:/p: namespaces and return the root element."""
    wrapped = (
        f'<root xmlns:a="{A_NS}" xmlns:p="{P_NS}" xmlns:r="{NSMAP["r"]}">'
        + fragment + "</root>"
    )
    return etree.fromstring(wrapped)[0]


# --------------------------------------------------------------------------- #
# Slide-management helpers
# --------------------------------------------------------------------------- #
def remove_all_slides(prs):
    """Drop the 5 demo slides from the template so we can add our own."""
    sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    for sldId in list(sldIdLst):
        rId = sldId.get(qn("r:id"))
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)


def _strip_layout_prompts(prs):
    """Strip 'Click to edit Master text styles' + 'March 2025' / 'Presenter Name'
    sample content from layout placeholders so unused placeholders render blank."""
    PROMPT_MARKERS = (
        "Click to edit", "Click here to add",
        "March 2025", "Presenter Name", "Lorem ipsum",
    )
    for layout in prs.slide_layouts:
        for shp in layout.shapes:
            if not shp.has_text_frame:
                continue
            tf_elem = shp.text_frame._txBody  # type: ignore[attr-defined]
            for r in list(tf_elem.iter(f"{{{A_NS}}}r")):
                t = r.find(f"{{{A_NS}}}t")
                if t is not None and t.text and any(m in t.text for m in PROMPT_MARKERS):
                    t.text = ""


def _delete_all_placeholders(slide):
    """Delete every placeholder from a content slide — we'll inject our own
    chrome (chip, title, body border, sldnum) per slides.md."""
    for shp in list(slide.placeholders):
        sp = shp._element  # type: ignore[attr-defined]
        sp.getparent().remove(sp)


def _sptree(slide):
    return slide.shapes._spTree  # type: ignore[attr-defined]


def _next_id(slide, start=10):
    """Get a fresh shape id that doesn't collide with existing shapes."""
    used = {int(e.get("id")) for e in _sptree(slide).iter(f"{{{P_NS}}}cNvPr") if e.get("id")}
    i = start
    while i in used:
        i += 1
    return i


# --------------------------------------------------------------------------- #
# CHROME — exact-XML injections per slides.md
# --------------------------------------------------------------------------- #
def add_body_border(slide):
    """Inject the body-border rectangle (REQUIRED). Added FIRST so it sits at
    the bottom of z-order, behind content."""
    sid = _next_id(slide)
    frag = f'''
<p:sp>
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="BodyBorder"/>
    <p:cNvSpPr/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="{BB_X}" y="{BB_Y}"/>
      <a:ext cx="{BB_W}" cy="{BB_H}"/>
    </a:xfrm>
    <a:prstGeom prst="roundRect">
      <a:avLst><a:gd name="adj" fmla="val 2185"/></a:avLst>
    </a:prstGeom>
    <a:solidFill><a:srgbClr val="FFFFFF"/></a:solidFill>
    <a:ln w="14605">
      <a:solidFill><a:srgbClr val="252525"/></a:solidFill>
      <a:prstDash val="solid"/>
    </a:ln>
  </p:spPr>
  <p:txBody><a:bodyPr/><a:lstStyle/><a:p><a:endParaRPr lang="en-US" sz="1350"/></a:p></p:txBody>
</p:sp>
'''
    _sptree(slide).append(_xml(frag))


def add_chip(slide, text=CHIP):
    """Inject the canonical pill chip (top-left, noFill, border #7300BF, adj=14493)."""
    sid = _next_id(slide)
    frag = f'''
<p:sp>
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="Chip"/>
    <p:cNvSpPr txBox="1"><a:spLocks/></p:cNvSpPr>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="196514" y="211211"/>
      <a:ext cx="896937" cy="304800"/>
    </a:xfrm>
    <a:prstGeom prst="roundRect">
      <a:avLst><a:gd name="adj" fmla="val 14493"/></a:avLst>
    </a:prstGeom>
    <a:noFill/>
    <a:ln>
      <a:solidFill><a:srgbClr val="7300BF"/></a:solidFill>
    </a:ln>
  </p:spPr>
  <p:txBody>
    <a:bodyPr anchor="ctr"><a:noAutofit/></a:bodyPr>
    <a:lstStyle/>
    <a:p>
      <a:pPr algn="ctr"/>
      <a:r>
        <a:rPr lang="en-US" sz="1000" b="1">
          <a:solidFill><a:srgbClr val="7300BF"/></a:solidFill>
          <a:latin typeface="Poppins"/>
        </a:rPr>
        <a:t>{text}</a:t>
      </a:r>
    </a:p>
  </p:txBody>
</p:sp>
'''
    _sptree(slide).append(_xml(frag))


def add_topic_title(slide, text):
    """Inject the topic title as a PLAIN TEXT BOX (NOT a p:ph placeholder).
    Y must center to the chip's center (~363611 EMU)."""
    sid = _next_id(slide)
    # Escape XML-special chars
    text_esc = text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
    frag = f'''
<p:sp>
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="Topic"/>
    <p:cNvSpPr txBox="1"/>
    <p:nvPr/>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="1210000" y="225000"/>
      <a:ext cx="6900000" cy="276999"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
    <a:noFill/>
  </p:spPr>
  <p:txBody>
    <a:bodyPr wrap="square" rtlCol="0"><a:spAutoFit/></a:bodyPr>
    <a:lstStyle/>
    <a:p>
      <a:r>
        <a:rPr lang="en-US" sz="1200" b="1" dirty="0">
          <a:latin typeface="Poppins"/>
        </a:rPr>
        <a:t>{text_esc}</a:t>
      </a:r>
    </a:p>
  </p:txBody>
</p:sp>
'''
    _sptree(slide).append(_xml(frag))


def add_sldnum(slide):
    """Inject the slide-number placeholder per slide so the page number renders."""
    sid = _next_id(slide)
    frag = f'''
<p:sp>
  <p:nvSpPr>
    <p:cNvPr id="{sid}" name="SlideNum"/>
    <p:cNvSpPr txBox="1"><a:spLocks noGrp="1"/></p:cNvSpPr>
    <p:nvPr><p:ph type="sldNum" idx="4"/></p:nvPr>
  </p:nvSpPr>
  <p:spPr>
    <a:xfrm>
      <a:off x="8213316" y="4847300"/>
      <a:ext cx="571800" cy="232800"/>
    </a:xfrm>
    <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
  </p:spPr>
  <p:txBody>
    <a:bodyPr spcFirstLastPara="1" wrap="square" lIns="91425" tIns="91425" rIns="91425" bIns="91425" anchor="ctr" anchorCtr="0">
      <a:noAutofit/>
    </a:bodyPr>
    <a:lstStyle/>
    <a:p><a:fld type="slidenum" id="{{00000000-1234-1234-1234-{sid:012x}}}"><a:rPr lang="en-US" sz="700" dirty="0"><a:latin typeface="Poppins"/></a:rPr></a:fld></a:p>
  </p:txBody>
</p:sp>
'''
    _sptree(slide).append(_xml(frag))


# --------------------------------------------------------------------------- #
# Content helpers (use python-pptx convenience for tables, images, text)
# --------------------------------------------------------------------------- #
def set_text(shape, lines, font_size=None, bold=False, color=INK, align=None,
             italic=False):
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(lines if isinstance(lines, list) else [lines]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if align is not None:
            p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = FONT
        if font_size is not None:
            r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.italic = italic
        if color is not None:
            r.font.color.rgb = color


def add_textbox(slide, x, y, w, h):
    return slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))


def add_bullets(slide, items, x, y, w, h, fs=11.5, sp=8):
    """Bullet list inside a textbox. Items: list of {lead, body}.
    Purple bullet+lead, black body. Bullet glyph drawn inline as '• '."""
    box = add_textbox(slide, x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if i > 0:
            p.space_before = Pt(sp)
        # Force-no list bullet glyph so we control bullets inline.
        pPr = p._p.get_or_add_pPr()
        for child in list(pPr):
            if child.tag in (f"{{{A_NS}}}buChar", f"{{{A_NS}}}buAutoNum", f"{{{A_NS}}}buFont"):
                pPr.remove(child)
        pPr.append(etree.SubElement(pPr, f"{{{A_NS}}}buNone"))
        pPr.set("marL", "0")
        pPr.set("indent", "0")
        r1 = p.add_run()
        r1.text = "•  " + it["lead"]
        r1.font.name = FONT
        r1.font.size = Pt(fs)
        r1.font.bold = True
        r1.font.color.rgb = PURPLE
        r2 = p.add_run()
        r2.text = it["body"]
        r2.font.name = FONT
        r2.font.size = Pt(fs)
        r2.font.color.rgb = INK


def add_image(slide, path, x, y, w, h):
    return slide.shapes.add_picture(path, Emu(x), Emu(y), Emu(w), Emu(h))


def add_table(slide, rows, x, y, w, h, col_widths_in=None, header_fill=DARKP,
              header_color=WHITE, body_color=INK, font_size=10,
              cell_padding_in=0.05):
    n_rows = len(rows)
    n_cols = len(rows[0])
    table = slide.shapes.add_table(n_rows, n_cols, Emu(x), Emu(y), Emu(w), Emu(h)).table
    if col_widths_in is not None:
        for ci, cw in enumerate(col_widths_in):
            table.columns[ci].width = Emu(int(Inches(cw)))
    for ri, row in enumerate(rows):
        is_header = (ri == 0)
        for ci, cell in enumerate(row):
            tc = table.cell(ri, ci)
            tc.margin_left = tc.margin_right = Emu(int(Inches(cell_padding_in)))
            tc.margin_top = tc.margin_bottom = Emu(int(Inches(0.03)))
            if isinstance(cell, dict):
                txt = cell.get("text", "")
                bold = cell.get("bold", False) or is_header
                color = cell.get("color", header_color if is_header else body_color)
                align = cell.get("align", None)
                italic = cell.get("italic", False)
            else:
                txt = cell
                bold = is_header
                color = header_color if is_header else body_color
                align = None
                italic = False
            if is_header:
                tc.fill.solid()
                tc.fill.fore_color.rgb = header_fill
            tf = tc.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            if align == "right":
                p.alignment = PP_ALIGN.RIGHT
            elif align == "center":
                p.alignment = PP_ALIGN.CENTER
            else:
                p.alignment = PP_ALIGN.LEFT
            r = p.add_run()
            r.text = str(txt)
            r.font.name = FONT
            r.font.size = Pt(font_size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.color.rgb = color
    return table


# --------------------------------------------------------------------------- #
# Slide builders
# --------------------------------------------------------------------------- #
def content_slide(prs, topic_title, build_body):
    """Standard content slide: white bg + body border + chip + title + sldnum +
    body content. Uses layout 2 (Blank White Layout) but deletes its
    placeholders so we own the chrome."""
    s = prs.slides.add_slide(prs.slide_layouts[2])
    _delete_all_placeholders(s)
    add_body_border(s)          # FIRST in z-order (sits behind content)
    add_chip(s, CHIP)
    add_topic_title(s, topic_title)
    add_sldnum(s)
    build_body(s)
    return s


def title_slide(prs):
    """Use the template's title-slide layout directly (its own gradient)."""
    s = prs.slides.add_slide(prs.slide_layouts[0])
    for shp in s.placeholders:
        if "Title" in shp.name:
            set_text(shp, ["MAMMAL", "Follow-up since 6/5"],
                     font_size=26, bold=True, color=WHITE)
        elif "Subtitle" in shp.name:
            set_text(shp,
                     ["Data gap vs model limit",
                      "Alternative models",
                      "The diversity finding",
                      "",
                      "Rohan Aryagondi  ·  June 2026"],
                     font_size=12, color=WHITE)
    return s


def closing_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[42])
    for shp in s.placeholders:
        if "Subtitle" in shp.name:
            set_text(shp,
                     ["docs/meeting_followup_report.md  ·  full detail in results/"],
                     font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
    return s


# --------------------------------------------------------------------------- #
# Per-slide content builders (each takes one `slide` arg and fills the body
# rectangle area: CONTENT_X / CONTENT_Y / CONTENT_W / CONTENT_H).
# --------------------------------------------------------------------------- #
def c_asks(s):
    items = [
        {"lead": "Diagnostic (Graham, Matt)", "body": " — is MAMMAL's Nav binding failure a data gap or a model limit? Decides whether a Quiver fine-tune helps."},
        {"lead": "Slide-9 alternative models (Matt, David)", "body": " — benchmark ConPLex / Boltz-2 / ADMET-AI / ESM-C on our data, not paper badges."},
        {"lead": "PROTON evaluation (Matt, David)", "body": " — Zitnik Lab's neurology relational FM, closer to the Sapphire KG use case than to binding prediction."},
        {"lead": "BBBP characterization (Graham, Mahdi)", "body": " — why is BBBP \"too permissive\"? Mahdi's reframe: trust the no's, investigate the yes's."},
    ]
    add_bullets(s, items, CONTENT_X, CONTENT_Y, CONTENT_W, CONTENT_H, fs=12, sp=10)


def c_tldr(s):
    items = [
        {"lead": "Data volume is necessary, not sufficient", "body": " — a minority of well-trained targets work, most don't. The high-data regime is bimodal."},
        {"lead": "No off-the-shelf alternative beats MAMMAL", "body": " — ConPLex sits at chance too; the failure is general to the BindingDB-trained DTI tooling space."},
        {"lead": "ADMET-AI finally fills the toxicity-gate slot", "body": " — DILI TPR 0.83 vs MAMMAL ClinTox 0.08. Deprecate ClinTox."},
        {"lead": "Protein-embedding pitch is at parity, not superiority", "body": " — MAMMAL ties open MIT-licensed ESM-2-650M on the canonical CRISPR-N panel."},
        {"lead": "mTOR is the next BRAF", "body": " — truncation falsified; chance even with kinase domain fully visible. Drop from any pitch."},
        {"lead": "Project thesis unchanged but sharper", "body": " — MAMMAL is commodity enrichment. Moat stays Quiver functional trace + V1-T."},
    ]
    add_bullets(s, items, CONTENT_X, CONTENT_Y, CONTENT_W, CONTENT_H, fs=11.5, sp=7)


def c_ceiling(s):
    rows = [
        ["Target", "Class", {"text": "n_pairs", "align": "right"},
         {"text": "AUROC random", "align": "right"},
         {"text": "AUROC matched", "align": "right"},
         {"text": "off-target Δ", "align": "right"}],
        [{"text": "RORC", "bold": True}, "nuclear receptor", {"text": "374", "align": "right"},
         {"text": "0.97", "bold": True, "color": PURPLE, "align": "right"},
         {"text": "0.95", "bold": True, "color": PURPLE, "align": "right"},
         {"text": "+0.68", "align": "right"}],
        [{"text": "CA2", "bold": True}, "other (CA)", {"text": "269", "align": "right"},
         {"text": "0.87", "bold": True, "color": PURPLE, "align": "right"},
         {"text": "0.84", "bold": True, "color": PURPLE, "align": "right"},
         {"text": "+1.97", "bold": True, "align": "right"}],
        [{"text": "Adrb2", "bold": True}, "GPCR", {"text": "211", "align": "right"},
         {"text": "0.87", "bold": True, "color": PURPLE, "align": "right"},
         {"text": "0.88", "bold": True, "color": PURPLE, "align": "right"},
         {"text": "+0.83", "align": "right"}],
        [{"text": "mTOR", "bold": True}, {"text": "kinase (Quiver)", "italic": True},
         {"text": "192", "align": "right"},
         {"text": "0.76", "align": "right"},
         {"text": "0.56", "color": RED, "align": "right"},
         {"text": "−1.12", "color": RED, "bold": True, "align": "right"}],
        [{"text": "BRAF", "bold": True}, {"text": "kinase (most-trained!)", "italic": True},
         {"text": "532", "align": "right"},
         {"text": "0.47", "color": RED, "align": "right"},
         {"text": "0.46", "color": RED, "align": "right"},
         {"text": "+1.18", "align": "right"}],
        [{"text": "HRH1", "bold": True}, "GPCR", {"text": "184", "align": "right"},
         {"text": "0.40", "color": RED, "align": "right"},
         {"text": "0.33", "color": RED, "align": "right"},
         {"text": "+0.68", "align": "right"}],
    ]
    table_h = int(Inches(2.6))
    add_table(s, rows, CONTENT_X, CONTENT_Y, CONTENT_W, table_h,
              col_widths_in=[0.85, 1.9, 0.8, 1.4, 1.4, 1.3], font_size=10)
    cap = add_textbox(s, CONTENT_X, CONTENT_Y + table_h + int(Inches(0.15)),
                      CONTENT_W, int(Inches(0.8)))
    cap.text_frame.word_wrap = True
    p = cap.text_frame.paragraphs[0]
    for txt, bold in [
        ("3 of 6 land at AUROC ≥ 0.80 on both random and MW-matched decoys. ", False),
        ("BRAF — the most-trained target in BindingDB — is at chance. mTOR collapses on matched and inverts off-target.", True),
    ]:
        r = p.add_run(); r.text = txt; r.font.name = FONT; r.font.size = Pt(10.5)
        r.font.color.rgb = INK; r.font.bold = bold


def c_curve(s):
    img_w_in = 5.2
    img_h_in = 3.3
    img_x = CONTENT_X + (CONTENT_W - int(Inches(img_w_in))) // 2
    add_image(s, FIG_CURVE, img_x, CONTENT_Y, int(Inches(img_w_in)), int(Inches(img_h_in)))
    cap = add_textbox(s, CONTENT_X, CONTENT_Y + int(Inches(img_h_in + 0.15)),
                      CONTENT_W, int(Inches(0.6)))
    cap.text_frame.word_wrap = True
    p = cap.text_frame.paragraphs[0]
    for txt, color, bold in [
        ("Non-monotonic. ", PURPLE, True),
        ("Peaks at the 40–149 bin (0.77 ± 0.09) and ", INK, False),
        ("drops back ", RED, True),
        ("to 0.60 in the 150+ bin with σ doubling. Data volume alone doesn't predict success at the high end.", INK, False),
    ]:
        r = p.add_run(); r.text = txt; r.font.name = FONT; r.font.size = Pt(10.5)
        r.font.color.rgb = color; r.font.bold = bold


def c_bimodality(s):
    items = [
        {"lead": "mTOR truncation theory — FALSIFIED", "body": " — kinase-domain window (451 aa, fully visible, no truncation) keeps AUROC at 0.50. mTOR is the next BRAF, not a truncation artifact."},
        {"lead": "Chemodiversity vs AUROC — Spearman = −0.83", "body": " — narrower binder sets correlate with higher AUROC across the 6 ceiling targets."},
        {"lead": "Memorisation hypothesis — REFUTED", "body": " — held-out-scaffold test: RORC AUROC 0.93 on the 94 % of binders that aren't the dominant scaffold (vs 0.97 in-scaffold). Ceiling wins are real generalisation."},
        {"lead": "Decoy-distance + output-variance — also REFUTED", "body": " — Spearman = +0.03 / +0.26. Three hypotheses tested; mechanism behind ρ = −0.83 remains genuinely open."},
    ]
    add_bullets(s, items, CONTENT_X, CONTENT_Y, CONTENT_W, CONTENT_H, fs=11.5, sp=8)


def c_conplex(s):
    rows = [
        ["Test",
         {"text": "MAMMAL", "align": "center"},
         {"text": "ConPLex", "align": "center"}],
        ["Correlation (Spearman vs pChEMBL, n=10)",
         {"text": "+0.43 PASS", "bold": True, "color": PURPLE, "align": "center"},
         {"text": "−0.03 FAIL", "color": RED, "align": "center"}],
        ["Named test (suzetrigine → Nav1.8, z-margin)",
         {"text": "−0.69 FAIL", "color": RED, "align": "center"},
         {"text": "−2.35 FAIL", "color": RED, "align": "center"}],
        ["Nav1.8 binder-vs-decoy AUROC",
         {"text": "0.43", "color": RED, "align": "center"},
         {"text": "0.39", "color": RED, "align": "center"}],
        ["mTOR binder-vs-decoy AUROC",
         {"text": "0.54", "color": RED, "align": "center"},
         {"text": "0.58", "color": RED, "align": "center"}],
        ["WDR91 SPR (fine-tuned MAMMAL vs zero-shot ConPLex)",
         {"text": "0.82 STRONG", "bold": True, "color": PURPLE, "align": "center"},
         {"text": "0.59", "align": "center"}],
    ]
    table_h = int(Inches(2.4))
    add_table(s, rows, CONTENT_X, CONTENT_Y, CONTENT_W, table_h,
              col_widths_in=[4.4, 2.0, 2.0], font_size=10.5)
    cap = add_textbox(s, CONTENT_X, CONTENT_Y + table_h + int(Inches(0.2)),
                      CONTENT_W, int(Inches(0.8)))
    cap.text_frame.word_wrap = True
    p = cap.text_frame.paragraphs[0]
    for txt, color, bold in [
        ("ConPLex does not beat MAMMAL anywhere. ", PURPLE, True),
        ("Both sit at chance on Nav1.8 / mTOR. ", INK, False),
        ("Zero-shot DTI failure is general to BindingDB-trained tooling — not MAMMAL-specific.", INK, True),
    ]:
        r = p.add_run(); r.text = txt; r.font.name = FONT; r.font.size = Pt(11)
        r.font.color.rgb = color; r.font.bold = bold


def c_admet(s):
    rows = [
        ["Endpoint",
         {"text": "AUROC", "align": "center"},
         {"text": "TPR (toxic)", "align": "center"},
         {"text": "TNR (safe)", "align": "center"}],
        [{"text": "ADMET-AI DILI", "bold": True},
         {"text": "0.73", "bold": True, "color": PURPLE, "align": "center"},
         {"text": "0.83", "bold": True, "color": PURPLE, "align": "center"},
         {"text": "0.67", "align": "center"}],
        ["ADMET-AI AMES",
         {"text": "0.67", "align": "center"},
         {"text": "0.17", "align": "center"},
         {"text": "0.93", "align": "center"}],
        ["ADMET-AI ClinTox",
         {"text": "0.50", "color": RED, "align": "center"},
         {"text": "0.00", "color": RED, "align": "center"},
         {"text": "1.00", "align": "center"}],
        ["ADMET-AI hERG",
         {"text": "0.48", "color": RED, "align": "center"},
         {"text": "0.33", "align": "center"},
         {"text": "0.40", "align": "center"}],
        [{"text": "MAMMAL ClinTox", "italic": True},
         {"text": "0.28", "color": RED, "align": "center"},
         {"text": "0.08", "color": RED, "align": "center"},
         {"text": "1.00", "align": "center"}],
    ]
    table_h = int(Inches(2.1))
    add_table(s, rows, CONTENT_X, CONTENT_Y, CONTENT_W, table_h,
              col_widths_in=[3.5, 1.8, 1.8, 1.8], font_size=10.5)
    items = [
        {"lead": "ADMET-AI DILI catches 10/12 toxics", "body": " — vs MAMMAL ClinTox's 1/12. Earns the toxicity-gate slot."},
        {"lead": "ClinTox dataset is the wrong endpoint", "body": " — both models fail. Use mechanism-specific endpoints (DILI / hERG / AMES). Deprecate MAMMAL ClinTox."},
    ]
    add_bullets(s, items, CONTENT_X, CONTENT_Y + table_h + int(Inches(0.25)),
                CONTENT_W, int(Inches(1.2)), fs=10.5, sp=6)


def c_esm(s):
    rows = [
        ["Model",
         {"text": "License", "align": "center"},
         {"text": "NN-recall", "align": "center"},
         {"text": "family gap", "align": "center"}],
        [{"text": "MAMMAL 458M", "bold": True}, {"text": "IBM (research)", "align": "center"},
         {"text": "0.750", "bold": True, "color": PURPLE, "align": "center"},
         {"text": "0.374", "align": "center"}],
        ["ESM-2 650M (raw cosine)", {"text": "MIT (open)", "align": "center"},
         {"text": "0.725", "align": "center"},
         {"text": "0.039", "align": "center"}],
        [{"text": "ESM-2 650M (centered — the textbook fix)", "bold": True},
         {"text": "MIT (open)", "align": "center"},
         {"text": "0.750 (tie)", "bold": True, "color": PURPLE, "align": "center"},
         {"text": "0.417", "bold": True, "align": "center"}],
    ]
    table_h = int(Inches(1.7))
    add_table(s, rows, CONTENT_X, CONTENT_Y, CONTENT_W, table_h,
              col_widths_in=[3.7, 2.5, 1.4, 1.4], font_size=10.5)
    cap = add_textbox(s, CONTENT_X, CONTENT_Y + table_h + int(Inches(0.05)),
                      CONTENT_W, int(Inches(0.3)))
    pp = cap.text_frame.paragraphs[0]
    r = pp.add_run(); r.text = "40-gene CRISPR-N panel  ·  same protocol  ·  mean-pool + cosine NN"
    r.font.name = FONT; r.font.size = Pt(10); r.font.italic = True; r.font.color.rgb = GREY
    items = [
        {"lead": "Prior 0.92 vs 0.88 was on a 25-protein toy panel", "body": " — never re-run on the canonical CRISPR-N until now. On the real panel under the standard anisotropy correction, MAMMAL ties."},
        {"lead": "ESM-2-650M is MIT-licensed and a clean swap", "body": " — Sapphire embedding-layer pitch survives at parity, not superiority. Open option on the table."},
    ]
    add_bullets(s, items, CONTENT_X, CONTENT_Y + table_h + int(Inches(0.45)),
                CONTENT_W, int(Inches(1.4)), fs=10.5, sp=6)


def c_bbbp(s):
    # Left: BBBP scatter PNG; Right: bullets + Mahdi quote
    LEFT_W = int(Inches(4.2))
    RIGHT_X = CONTENT_X + LEFT_W + int(Inches(0.2))
    RIGHT_W = CONTENT_W - LEFT_W - int(Inches(0.2))
    add_image(s, FIG_BBBP, CONTENT_X, CONTENT_Y, LEFT_W, int(Inches(3.5)))
    items = [
        {"lead": "Spearman(P(BBB+), MW) = −0.73", "body": " — Graham directionally right. HBA −0.67, TPSA −0.61."},
        {"lead": "Size + polarity exclusion gate", "body": " — operative cliff is ≳450 Da + polar → exclude. Not \"<300 Da → brain\"."},
        {"lead": "8/8 \"predicted no\" are truly non-penetrant", "body": " — TNR perfect on the <0.3 set."},
        {"lead": "29/43 \"predicted yes\" are CNS-active (67 %)", "body": " — only marginally above base rate."},
        {"lead": "Saturated 0/1", "body": " — no drug lands between 0.3 and 0.7."},
    ]
    add_bullets(s, items, RIGHT_X, CONTENT_Y, RIGHT_W, int(Inches(2.8)), fs=10.5, sp=6)
    cap = add_textbox(s, RIGHT_X, CONTENT_Y + int(Inches(2.95)),
                      RIGHT_W, int(Inches(0.4)))
    pp = cap.text_frame.paragraphs[0]
    r = pp.add_run()
    r.text = "Mahdi's reframe — \"trust the no's, investigate the yes's\" — fully validated."
    r.font.name = FONT; r.font.size = Pt(10.5); r.font.bold = True; r.font.color.rgb = PURPLE


def c_aws_status(s):
    items = [
        {"lead": "PROTON access — CONFIRMED open", "body": " — code (mims-harvard/PROTON, MIT) + weights (HuggingFace) + NeuroKG (Harvard Dataverse). No licensing friction."},
        {"lead": "PROTON install on macOS ARM64 — impractical", "body": " — DGL ships no Apple-Silicon wheels; from-source compile ~1 hr, disk at 96 %. Eval needs Linux or AWS g5."},
        {"lead": "Boltz-2 — pilot launched then cancelled at $0.11", "body": " — g5.xlarge spun up, volume attached, killed on direction. Both Boltz-2 and PROTON queue for the same AWS g5 when AWS is back on."},
        {"lead": "Planned PROTON eval (when AWS is back)", "body": " — (i) embed CRISPR-N panel, compare NN-recall vs MAMMAL/ESM-2; (ii) drug-target link prediction for Nav1.8/mTOR; (iii) hypothesis generation on one Quiver CRISPR-N hit."},
    ]
    add_bullets(s, items, CONTENT_X, CONTENT_Y, CONTENT_W, CONTENT_H, fs=11, sp=8)


def c_changed(s):
    rows = [
        ["Prior belief", "Updated reading"],
        ["Nav fails because BindingDB has no Nav → close the gap = done.",
         {"text": "True but insufficient. The high end is bimodal; data volume alone doesn't predict success.", "bold": True}],
        ["Off-the-shelf alternatives might rescue Nav (ConPLex).",
         {"text": "No off-the-shelf alternative beats MAMMAL on our targets. Boltz-2 is the only remaining unknown.", "bold": True}],
        ["MAMMAL is data-suited on rich targets like mTOR.",
         {"text": "mTOR is the next BRAF. Truncation falsified. Drop from any pitch.", "bold": True, "color": RED}],
        ["MAMMAL ClinTox covers tox de-risking.",
         {"text": "ClinTox dataset is the wrong endpoint. ADMET-AI DILI is the replacement (TPR 0.83).", "bold": True}],
        ["MAMMAL embeddings beat ESM-2 → Sapphire layer = MAMMAL.",
         {"text": "Parity, not superiority. ESM-2-650M (MIT, open) ties on the canonical panel.", "bold": True}],
        ["BBBP head is reliable.",
         {"text": "Asymmetric. \"No\" is perfect; \"yes\" marginally above base rate. Soft positive flag, hard negative gate.", "bold": True}],
    ]
    add_table(s, rows, CONTENT_X, CONTENT_Y, CONTENT_W, int(Inches(3.6)),
              col_widths_in=[3.7, 5.0], font_size=9.5, cell_padding_in=0.06)


# --------------------------------------------------------------------------- #
def main():
    prs = Presentation(str(TEMPLATE))
    _strip_layout_prompts(prs)
    remove_all_slides(prs)

    title_slide(prs)
    content_slide(prs, "What the 6/5 check-in asked", c_asks)
    content_slide(prs, "Bottom line", c_tldr)
    content_slide(prs, "Where is MAMMAL data-suited? — ceiling test", c_ceiling)
    content_slide(prs, "Threshold curve — AUROC vs training pairs", c_curve)
    content_slide(prs, "Why does it fail on some rich-data targets?", c_bimodality)
    content_slide(prs, "ConPLex — does it beat MAMMAL? No.", c_conplex)
    content_slide(prs, "ADMET-AI replaces the broken ClinTox head", c_admet)
    content_slide(prs, "Protein embeddings — at parity, not superiority", c_esm)
    content_slide(prs, "BBBP head — what is it actually doing?", c_bbbp)
    content_slide(prs, "PROTON + Boltz-2 — both need AWS, both on hold", c_aws_status)
    content_slide(prs, "What changed vs going-in beliefs", c_changed)
    closing_slide(prs)

    prs.save(str(OUTPUT))
    print(f"WROTE {OUTPUT}")


if __name__ == "__main__":
    main()
